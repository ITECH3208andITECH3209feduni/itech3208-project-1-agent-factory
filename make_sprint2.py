"""
Agent Factory — Sprint 2 Showcase Deck  (v4 — equal member slides)
9 slides: Title · 5 × Member · Delivery · Demo · Thank You
10" × 5.62", dark navy, gold/cyan accents, left gold bar.
"""

import os, subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0A, 0x0F, 0x26)
HERO_BG   = RGBColor(0x06, 0x09, 0x1A)
HDR_BG    = RGBColor(0x04, 0x07, 0x12)
CARD_D    = RGBColor(0x0F, 0x17, 0x38)
CARD_L    = RGBColor(0x15, 0x20, 0x48)
FOOTER_BG = RGBColor(0x03, 0x05, 0x10)
GOLD      = RGBColor(0xD4, 0xAF, 0x37)
CYAN      = RGBColor(0x00, 0xC2, 0xD4)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W     = RGBColor(0xD8, 0xE8, 0xFF)
MUTED     = RGBColor(0x55, 0x70, 0xA0)
TBLHDR    = RGBColor(0x20, 0x30, 0x70)
RIGHTCOL  = RGBColor(0x0C, 0x12, 0x2C)
LT_BLUE   = RGBColor(0x80, 0xCC, 0xFF)
AMBER     = RGBColor(0xFF, 0xC0, 0x00)

W = 10.0
H = 5.62

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ── Core drawing helpers ──────────────────────────────────────────

def new_slide(bg=DARK_BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(W), Inches(H))
    _solid(r, bg)
    return s

def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def R(s, l, t, w, h, color):
    shp = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    _solid(shp, color)
    return shp

def T(s, l, t, w, h, text, size=10, bold=False, color=WHITE,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text        = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = "Calibri"
    return tb

def T_multi(s, l, t, w, h, lines, size=10, default_color=OFF_W):
    """lines = list of (text, bold, color|None). Each line is a paragraph."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, bold, color in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or default_color
        r.font.name = "Calibri"

def footer(s, text="Agent Factory   ·   Sprint 2 Showcase   ·   Federation University   ·   2026"):
    R(s, 0, H - 0.18, W, 0.18, FOOTER_BG)
    T(s, 0.22, H - 0.17, 9.50, 0.16, text, size=7.5, color=MUTED)

def stat_card(s, l, t, w, h, num, label, accent=GOLD):
    R(s, l, t, w, h, CARD_D)
    R(s, l, t, w, 0.04, accent)
    T(s, l, t + 0.04, w, 0.82, num,
      size=50, bold=True, color=accent, align=PP_ALIGN.CENTER)
    T(s, l, t + 0.86, w, 0.42, label,
      size=9.5, color=MUTED, align=PP_ALIGN.CENTER)


# ── Member slide template (identical for all 5 members) ──────────

def member_slide(name, role, badge, accent, stat_num, stat_label, contributions):
    """
    contributions = list of exactly 3 tuples:
        (title, proj_tag, body_lines)
        body_lines = list of (text, bold, color)
    """
    s = new_slide()

    # ─ Header ─
    R(s, 0,    0,    0.14, H,    GOLD)          # left gold bar
    R(s, 0.14, 0,    9.86, 0.78, HDR_BG)        # header bg
    R(s, 0.14, 0.78, 9.86, 0.03, GOLD)          # gold underline
    T(s, 0.34, 0.12, 7.20, 0.58,
      f"Sprint 2  ·  {name}",
      size=22, bold=True, color=WHITE)
    R(s, 7.88, 0.17, 1.98, 0.44, accent)
    T(s, 7.88, 0.26, 1.98, 0.44, badge,
      size=9, bold=True, color=HERO_BG, align=PP_ALIGN.CENTER)

    # ─ Left identity card ─
    LW = 3.26
    R(s, 0.22, 0.94, LW, 4.50, CARD_D)
    R(s, 0.22, 0.94, LW, 0.05, accent)

    T(s, 0.36, 1.06, LW - 0.18, 0.68, name,
      size=22, bold=True, color=accent)
    T(s, 0.36, 1.78, LW - 0.18, 0.28, role,
      size=10.5, color=MUTED)
    R(s, 0.36, 2.12, LW - 0.28, 0.02, accent)

    # Big stat
    T(s, 0.22, 2.20, LW, 0.86, stat_num,
      size=58, bold=True, color=accent, align=PP_ALIGN.CENTER)
    T(s, 0.22, 3.06, LW, 0.36, stat_label,
      size=10, color=MUTED, align=PP_ALIGN.CENTER)

    R(s, 0.36, 3.50, LW - 0.28, 0.02, CARD_L)
    T(s, 0.22, 3.58, LW, 0.26,
      "Sprint 2  ·  Agent Factory",
      size=8.5, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    # ─ Right: 3 contribution cards ─
    RX = 0.22 + LW + 0.10   # 3.58"
    RW = 9.80 - RX           # 6.22"
    CH  = 1.43
    GAP = 0.08

    for i, (title, proj_tag, body_lines) in enumerate(contributions):
        ty = 0.94 + i * (CH + GAP)
        R(s, RX, ty, RW, CH, CARD_D)
        R(s, RX, ty, RW, 0.04, accent)
        T(s, RX + 0.12, ty + 0.10, 4.80, 0.30,
          title, size=11.5, bold=True, color=accent)
        if proj_tag:
            T(s, RX + RW - 1.14, ty + 0.10, 1.04, 0.26,
              proj_tag, size=7.5, color=MUTED)
        T_multi(s, RX + 0.12, ty + 0.46, RW - 0.20, CH - 0.56,
                body_lines, size=9.5)

    footer(s)
    return s


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HERO
# ════════════════════════════════════════════════════════════════
s = new_slide(HERO_BG)

R(s, 0,    0,    0.32, H,    GOLD)
R(s, 7.80, 0,    2.20, 0.08, GOLD)
R(s, 9.50, 0,    0.50, H,    RIGHTCOL)

R(s, 7.68, 0.18, 2.14, 0.50, GOLD)
T(s, 7.68, 0.32, 2.06, 0.20, "SPRINT 2   ·   COMPLETE",
  size=8.5, bold=True, color=HERO_BG)

T(s, 0.53, 0.52, 9.00, 0.20,
  "ITECH3208   ·   FEDERATION UNIVERSITY",
  size=8.5, color=GOLD)

T(s, 0.52, 0.68, 9.00, 1.10, "Agent",  size=72, bold=True, color=WHITE)
T(s, 0.52, 1.60, 9.00, 1.00, "Factory", size=72, bold=True, color=GOLD)

R(s, 0.52, 2.64, 4.80, 0.03, GOLD)

T(s, 0.52, 2.76, 7.50, 0.45,
  "Autonomous AI Research Intelligence — 24 / 7",
  size=16, color=OFF_W)

for i, (icon, label) in enumerate([
    ("📚", "Literature Research"),
    ("🛒", "Amazon Intelligence"),
    ("🌐", "Web UI & REST API"),
]):
    lx = 0.52 + i * 3.08
    R(s, lx, 3.36, 2.88, 0.40, CARD_L)
    R(s, lx, 3.36, 0.04, 0.40, GOLD)
    T(s, lx + 0.10, 3.38, 2.74, 0.36, f"{icon}  {label}",
      size=10.5, bold=True, color=OFF_W)

T(s, 0.52, 3.96, 8.80, 0.36,
  "Dilraj Singh  ·  Dhiman Roy  ·  Prabhjot Singh  ·  Sahil K C  ·  Saifur Rahman B",
  size=10.0, color=MUTED)

R(s, 0.52, 4.42, 8.80, 0.03, GOLD)
T(s, 0.52, 4.50, 8.80, 0.32,
  "Powered by Claude AI   ·   Playwright   ·   Semantic Scholar API   ·   arXiv   ·   FastAPI",
  size=9.0, color=MUTED)

R(s, 0, 5.30, W, 0.33, FOOTER_BG)
T(s, 0.52, 5.33, 9.00, 0.26,
  "Confidential — Sprint 2 Showcase   ·   Federation University   ·   2026",
  size=8.0, color=MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — DILRAJ SINGH
# ════════════════════════════════════════════════════════════════
member_slide(
    name="Dilraj Singh",
    role="Lead Developer",
    badge="LEAD DEV",
    accent=GOLD,
    stat_num="67",
    stat_label="Jira Tickets Closed",
    contributions=[
        ("Literature Research Skill", "PROJ-127–132", [
            ("Multi-source fetcher: arXiv + Semantic Scholar.", False, OFF_W),
            ("Claude Haiku synthesis — summary, gaps, citations.", False, OFF_W),
            ("PaperCard model, rate-limit retry, demo script.", False, OFF_W),
        ]),
        ("SQLite Session Memory", "PROJ-135–155", [
            ("Replaced flat JSON with full SQLite schema.", False, OFF_W),
            ("SessionMemory class — sessions, messages, metadata.", False, OFF_W),
            ("Agent reads/writes DB every turn — no data loss.", False, OFF_W),
        ]),
        ("UI Sprint — Attachments & Seller Tools", "PROJ-156–194", [
            ("File attachments: PDF/DOCX/image parsing in chatbox.", False, OFF_W),
            ("Literature AI: dual-panel (Research Search + Integrity).", False, OFF_W),
            ("PPC builder card fix · arXiv 45 s timeout · chip UX.", False, OFF_W),
        ]),
    ]
)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — DHIMAN ROY
# ════════════════════════════════════════════════════════════════
member_slide(
    name="Dhiman Roy",
    role="Amazon Intelligence Lead",
    badge="AMAZON LEAD",
    accent=CYAN,
    stat_num="31",
    stat_label="Jira Tickets Closed",
    contributions=[
        ("Amazon Skill Core", "PROJ-111–113", [
            ("amazon_skill.py — full skill pipeline built.", False, OFF_W),
            ("amazon_cards.py — ProductCard + formatters.", False, OFF_W),
            ("CLI demo script for standalone testing.", False, OFF_W),
        ]),
        ("Playwright + RapidAPI + AI Scorer", "PROJ-166–169", [
            ("Stealth Playwright scrape: price, rating, Prime.", False, OFF_W),
            ("RapidAPI fallback when scraper is blocked.", False, OFF_W),
            ("0–100 AI score formula + MD5 TTL cache layer.", False, OFF_W),
        ]),
        ("Amazon Route + Claude Prompts", "PROJ-170–174", [
            ("GET /api/amazon FastAPI endpoint.", False, OFF_W),
            ("SEARCH, COMPARE, RECOMMEND prompt templates.", False, OFF_W),
            ("Pydantic Product + AmazonResponse schemas.", False, OFF_W),
        ]),
    ]
)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — PRABHJOT SINGH
# ════════════════════════════════════════════════════════════════
member_slide(
    name="Prabhjot Singh",
    role="Infrastructure & DevOps",
    badge="INFRA",
    accent=GOLD,
    stat_num="27",
    stat_label="Jira Tickets Closed",
    contributions=[
        ("Docker & Container Environment", "", [
            ("Maintained docker-compose.yml across the sprint.", False, OFF_W),
            ("Service health checks and restart policies.", False, OFF_W),
            ("BuildKit cache tuning — eliminated stale layers.", False, OFF_W),
        ]),
        ("CI/CD & Code Review", "", [
            ("Reviewed PRs across Amazon & Literature skills.", False, OFF_W),
            ("Caught integration conflicts before they merged.", False, OFF_W),
            ("Kept the main branch green throughout the sprint.", False, OFF_W),
        ]),
        ("Environment & Dev Setup", "", [
            (".env management — centralised secret config.", False, OFF_W),
            ("Port/path config for Uvicorn, SQLite, Playwright.", False, OFF_W),
            ("Dev onboarding guide — one-command setup.", False, OFF_W),
        ]),
    ]
)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — SAHIL K C
# ════════════════════════════════════════════════════════════════
member_slide(
    name="Sahil K C",
    role="QA & Documentation",
    badge="QA & DOCS",
    accent=CYAN,
    stat_num="27",
    stat_label="Jira Tickets Closed",
    contributions=[
        ("End-to-End Feature Testing", "", [
            ("Verified Literature: arXiv + Semantic Scholar.", False, OFF_W),
            ("Confirmed Amazon cards: scores, Prime, links.", False, OFF_W),
            ("Tested REST API responses against contracts.", False, OFF_W),
        ]),
        ("Web UI Acceptance Testing", "", [
            ("Tested localhost:8000 and /literature pages.", False, OFF_W),
            ("Verified error banners and empty-state messages.", False, OFF_W),
            ("Confirmed rate-limit notices display correctly.", False, OFF_W),
        ]),
        ("Sprint Documentation", "", [
            ("Maintained sprint notes and meeting records.", False, OFF_W),
            ("Reviewed API usage docs for accuracy.", False, OFF_W),
            ("Error message copy review and final sign-off.", False, OFF_W),
        ]),
    ]
)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — SAIFUR RAHMAN BHUIYAN
# ════════════════════════════════════════════════════════════════
member_slide(
    name="Saifur Rahman B",
    role="QA Lead",
    badge="QA LEAD",
    accent=GOLD,
    stat_num="24",
    stat_label="Jira Tickets Closed",
    contributions=[
        ("Acceptance Criteria & Test Design", "", [
            ("Wrote acceptance criteria for all new skills.", False, OFF_W),
            ("Test cases designed before dev started.", False, OFF_W),
            ("Covered Literature, Amazon, API, and Web UI.", False, OFF_W),
        ]),
        ("Pydantic Schema Verification", "", [
            ("Reviewed Paper + LiteratureResponse schemas.", False, OFF_W),
            ("Reviewed Product + AmazonResponse schemas.", False, OFF_W),
            ("Caught field-type mismatches before merge.", False, OFF_W),
        ]),
        ("Input / Output Contracts", "", [
            ("Defined expected inputs and outputs per skill.", False, OFF_W),
            ("Ensured Claude AI responses matched format.", False, OFF_W),
            ("Verified end-to-end output consistency.", False, OFF_W),
        ]),
    ]
)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — EVERY COMMITMENT DELIVERED
# ════════════════════════════════════════════════════════════════
s = new_slide()
R(s, 0,    0,    0.14, H,    GOLD)
R(s, 0.14, 0,    9.86, 0.68, HDR_BG)
R(s, 0.14, 0.68, 9.86, 0.03, GOLD)
T(s, 0.32, 0.09, 7.50, 0.52,
  "Sprint 2: Every Commitment Delivered",
  size=23, bold=True, color=WHITE)
R(s, 8.02, 0.15, 1.84, 0.38, GOLD)
T(s, 8.02, 0.24, 1.84, 0.38, "EXECUTION PROOF",
  size=8, bold=True, color=HERO_BG, align=PP_ALIGN.CENTER)

# Big stat row
for i, (num, lbl, ac) in enumerate([
    ("176",  "Jira Tickets\nClosed",  GOLD),
    ("14",   "GitHub PRs\nMerged",    CYAN),
    ("5",    "Team\nMembers",         GOLD),
    ("178",  "Total\nTickets",        CYAN),
]):
    stat_card(s, 0.22 + i * 2.41, 0.78, 2.33, 1.32, num, lbl, ac)

# Delivery table header
R(s, 0.22, 2.18, 9.52, 0.28, CARD_L)
for lx, w, txt in [
    (0.28, 3.68, "Epic / Workstream"),
    (4.06, 1.04, "Owner"),
    (5.20, 0.62, "Tickets"),
    (5.92, 2.62, "Deliverables"),
    (8.64, 1.22, "Status"),
]:
    T(s, lx, 2.23, w, 0.22, txt, size=8.5, bold=True, color=GOLD)

R(s, 0.22, 2.46, 9.52, 0.03, GOLD)

rows = [
    (CARD_D, GOLD, "Literature Skill — arXiv + Semantic Scholar + Synthesis",
     "Dilraj", "6",  "Fetcher · Claude synthesis · PaperCard · Demo",   "✅ Done"),
    (CARD_L, CYAN, "Amazon Skill — Scraper + Scorer + Cache + API",
     "Dhiman", "12", "Playwright · RapidAPI · 0–100 score · TTL cache", "✅ Done"),
    (CARD_D, GOLD, "REST API — /api/literature + /api/amazon + schemas",
     "Both",   "4",  "FastAPI routes · Pydantic schemas · Swagger docs",  "✅ Done"),
    (CARD_L, CYAN, "Web UI — Browser search interface (localhost:8000)",
     "Dilraj", "4",  "FastAPI static · literature.html · index.html",    "✅ Done"),
    (CARD_D, GOLD, "Session Memory — SQLite persistence across restarts",
     "Dilraj", "5",  "Schema · SessionMemory class · agent integration", "✅ Done"),
    (CARD_L, CYAN, "UI Polish — PPC cards · Literature dual-panel · chip UX · arXiv fix",
     "Dilraj", "8",  "Mode-type routing · inline S2 warning · query cleaner", "✅ Done"),
    (CARD_D, GOLD, "File Attachment System — PDF · DOCX · images in chatbox",
     "Dilraj", "3",  "PDF.js + mammoth.js lazy-load · preview chips · violet UX", "✅ Done"),
]

ROW_H  = 0.41   # tighter pitch — fits 7 rows before footer
ROW_TY = 0.03   # text y-offset within row

for i, (bg, ac, epic, owner, tickets, deliverables, status) in enumerate(rows):
    ry = 2.48 + i * ROW_H
    ty = ry  + ROW_TY
    R(s, 0.22, ry, 9.52, ROW_H, bg)
    R(s, 0.22, ry, 0.04, ROW_H, ac)
    T(s, 0.30, ty, 3.64, 0.34, epic,         size=8.0, color=OFF_W)
    T(s, 4.06, ty, 1.00, 0.34, owner,        size=8.0, color=MUTED)
    T(s, 5.20, ty, 0.58, 0.34, tickets,      size=8.0, color=ac, bold=True)
    T(s, 5.92, ty, 2.62, 0.34, deliverables, size=7.5, color=MUTED)
    T(s, 8.64, ty, 1.20, 0.34, status,       size=8.0, color=GOLD, bold=True)

footer(s)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — LIVE DEMO
# ════════════════════════════════════════════════════════════════
s = new_slide()
R(s, 0,    0,    0.14, H,    GOLD)
R(s, 0.14, 0,    9.86, 0.78, HDR_BG)
R(s, 0.14, 0.78, 9.86, 0.03, GOLD)
T(s, 0.34, 0.10, 7.60, 0.58, "Live Demo", size=24, bold=True, color=WHITE)
R(s, 7.90, 0.17, 1.96, 0.44, CYAN)
T(s, 7.90, 0.26, 1.96, 0.44, "TRY IT NOW",
  size=9, bold=True, color=HERO_BG, align=PP_ALIGN.CENTER)

for i, (icon, title, ac, url, steps) in enumerate([
    ("📚", "Literature Search",     GOLD, "http://localhost:8000/literature",
     ["Type a research topic",
      "→ arXiv + Semantic Scholar results",
      "→ Claude AI synthesis paragraph",
      "→ Research gaps identified"]),
    ("🛒", "Amazon Product Search", CYAN, "http://localhost:8000",
     ["Type a product query",
      "→ Scraped & AI-scored products",
      "→ Price · rating · reviews · Prime",
      "→ Direct Amazon links"]),
    ("📎", "File Attachment System", GOLD, "Any chatbox — click 📎 to attach",
     ["Attach PDF, DOCX, image or TXT",
      "→ PDF.js extracts text client-side",
      "→ mammoth.js parses Word files",
      "→ Content auto-sent with query"]),
]):
    lx = 0.22 + i * 3.22
    R(s, lx, 0.94, 3.04, 4.52, CARD_D)
    R(s, lx, 0.94, 3.04, 0.05, ac)
    T(s, lx + 0.10, 1.04, 2.84, 0.36, f"{icon}  {title}",
      size=12.5, bold=True, color=ac)
    T(s, lx + 0.10, 1.46, 2.84, 0.22, url, size=9, color=MUTED, italic=True)
    for j, step in enumerate(steps):
        T(s, lx + 0.10, 1.78 + j * 0.66, 2.84, 0.60, step, size=10.5, color=OFF_W)

footer(s)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — SPRINT 3 ROADMAP  (AI Receptionist)
# ════════════════════════════════════════════════════════════════
s = new_slide()
R(s, 0,    0,    0.14, H,    GOLD)
R(s, 0.14, 0,    9.86, 0.78, HDR_BG)
R(s, 0.14, 0.78, 9.86, 0.03, GOLD)
T(s, 0.34, 0.10, 7.20, 0.58, "Sprint 3 — What's Next",
  size=24, bold=True, color=WHITE)
R(s, 7.88, 0.17, 1.98, 0.44, CYAN)
T(s, 7.88, 0.26, 1.98, 0.44, "ROADMAP",
  size=9, bold=True, color=HERO_BG, align=PP_ALIGN.CENTER)

# ── Hero card: AI Receptionist ──
R(s, 0.22, 0.94, 9.56, 1.00, CARD_D)
R(s, 0.22, 0.94, 9.56, 0.05, GOLD)
T(s, 0.40, 1.04, 1.20, 0.36, "🤖", size=22, color=GOLD)
T(s, 1.54, 1.02, 4.80, 0.34,
  "AI Receptionist — Automated Call Centre Agent",
  size=13, bold=True, color=GOLD)
T(s, 1.54, 1.38, 7.80, 0.46,
  "A 24/7 AI agent that replaces a human receptionist. Deployed on Mac Mini via OpenClaw — "
  "answers enquiries, books appointments, routes to the right agent, and escalates to humans.",
  size=9.5, color=OFF_W)

# ── 6 feature cards — 3 columns × 2 rows ──
features = [
    ("💬  Telegram Bot",        GOLD,
     "Talk to your agents from\nyour phone. No browser\nneeded. Live 24/7 via\nMac Mini OpenClaw."),
    ("🧠  Smart Intent Router", CYAN,
     "Claude reads your message\nand routes to the right\nagent automatically.\nNo menus or /commands."),
    ("📚  Knowledge Base",      GOLD,
     "Upload FAQs, docs,\nservices & pricing per\nbusiness client. Agent\nanswers from real data."),
    ("📅  Appointment Booking", CYAN,
     "Book via Google Calendar\ndirectly from the chat.\nNo human needed.\nInstant confirmation."),
    ("🚨  Escalation Engine",   GOLD,
     "Hands complex cases to\na human with full\nconversation summary.\nNever drops context."),
    ("⚙️  Multi-Client Setup",  CYAN,
     "Configure the agent for\nany business client\nindependently. One\nplatform, many clients."),
]

CW = 3.04
CH = 1.42
GAP_X = 0.11
GAP_Y = 0.08
START_X = 0.22
START_Y = 2.04

for i, (title, ac, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    lx = START_X + col * (CW + GAP_X)
    ty = START_Y + row * (CH + GAP_Y)
    R(s, lx, ty, CW, CH, CARD_D)
    R(s, lx, ty, CW, 0.04, ac)
    T(s, lx + 0.12, ty + 0.10, CW - 0.20, 0.30, title,
      size=10.5, bold=True, color=ac)
    T(s, lx + 0.12, ty + 0.44, CW - 0.20, CH - 0.52, desc,
      size=9, color=OFF_W)

footer(s)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — THANK YOU
# ════════════════════════════════════════════════════════════════
s = new_slide(HERO_BG)

R(s, 0,    0,    0.32, H,    GOLD)
R(s, 7.80, 0,    2.20, 0.08, GOLD)
R(s, 9.50, 0,    0.50, H,    RIGHTCOL)

R(s, 7.68, 0.18, 2.14, 0.50, GOLD)
T(s, 7.68, 0.32, 2.06, 0.20, "SPRINT 2   ·   COMPLETE",
  size=8.5, bold=True, color=HERO_BG)

T(s, 0.53, 0.52, 9.00, 0.20,
  "ITECH3208   ·   FEDERATION UNIVERSITY",
  size=8.5, color=GOLD)

T(s, 0.52, 0.80, 9.00, 1.10, "Thank",  size=72, bold=True, color=WHITE)
T(s, 0.52, 1.72, 9.00, 1.00, "You.",   size=72, bold=True, color=GOLD)

R(s, 0.52, 2.76, 4.80, 0.03, GOLD)

T(s, 0.52, 2.90, 7.50, 0.45, "Any questions?", size=20, color=OFF_W)

T(s, 0.52, 3.50, 8.80, 0.36,
  "Dilraj Singh  ·  Dhiman Roy  ·  Prabhjot Singh  ·  Sahil K C  ·  Saifur Rahman B",
  size=10.0, color=MUTED)

R(s, 0.52, 4.10, 8.80, 0.03, GOLD)
T(s, 0.52, 4.18, 8.80, 0.32,
  "github.com/ITECH3208andITECH3209feduni/itech3208-project-1-agent-factory",
  size=9.0, color=MUTED)

for i, (icon, label) in enumerate([
    ("📚", "Literature Research"),
    ("🛒", "Amazon Intelligence"),
    ("🌐", "Web UI & REST API"),
]):
    lx = 0.52 + i * 3.08
    R(s, lx, 4.60, 2.88, 0.40, CARD_L)
    R(s, lx, 4.60, 0.04, 0.40, GOLD)
    T(s, lx + 0.10, 4.62, 2.74, 0.36, f"{icon}  {label}",
      size=10.5, bold=True, color=OFF_W)

R(s, 0, 5.30, W, 0.33, FOOTER_BG)
T(s, 0.52, 5.33, 9.00, 0.26,
  "Confidential — Sprint 2 Showcase   ·   Federation University   ·   2026",
  size=8.0, color=MUTED)


# ── Save & open ──────────────────────────────────────────────────
OUT = os.path.expanduser("~/Desktop/Agent_Factory_Sprint2.pptx")
prs.save(OUT)
print(f"Saved → {OUT}")
subprocess.Popen(["open", OUT])
